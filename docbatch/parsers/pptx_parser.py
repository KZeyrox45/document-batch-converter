"""
PPTX document parser using python-pptx.

This module provides functionality to parse PowerPoint presentations
and extract text, tables, images, and metadata into structured JSON format.
"""

import re
import time
from pathlib import Path
from typing import List, Optional

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from docbatch.parsers.base import BaseParser
from docbatch.models import (
    DocumentOutput,
    DocumentMetadata,
    DocumentType,
    Section,
    SlideContent,
    TableData,
    ImageInfo,
)


class PPTXParser(BaseParser):
    """
    Parser for PPTX presentations using python-pptx.
    
    Extracts:
    - Slide text and speaker notes
    - Tables with headers and rows
    - Image metadata (dimensions)
    - Presentation metadata
    """
    
    SUPPORTED_EXTENSIONS = ['.pptx', '.ppt']
    
    def parse(self, filepath: str) -> DocumentOutput:
        """
        Parse a PPTX file and return structured output.
        
        Args:
            filepath: Path to the PPTX file.
            
        Returns:
            DocumentOutput containing parsed content.
        """
        self.clear_warnings()
        start_time = time.perf_counter()
        
        filename = Path(filepath).name
        self.logger.info(f"Parsing PPTX: {filename}")
        
        # Extract all components
        metadata = self.extract_metadata(filepath)
        slides = self.extract_slides(filepath)
        tables = self.extract_tables(filepath)
        images = self.extract_images(filepath)
        
        # Create sections from slides
        sections = self._create_sections_from_slides(slides)
        
        elapsed = time.perf_counter() - start_time
        
        return DocumentOutput(
            filename=filename,
            metadata=metadata,
            sections=sections,
            slides=slides,
            tables=tables,
            images=images,
            warnings=self.warnings.copy(),
            conversion_time=elapsed,
        )
    
    def extract_metadata(self, filepath: str) -> DocumentMetadata:
        """Extract metadata from PPTX file."""
        try:
            prs = Presentation(filepath)
            
            # Get presentation properties
            props = prs.core_properties
            
            return DocumentMetadata(
                file_type=DocumentType.PPTX.value,
                slides=len(prs.slides),
                author=props.author,
                title=props.title,
                subject=props.subject,
                creator=props.author,
                created=str(props.created) if props.created else None,
                modified=str(props.modified) if props.modified else None,
            )
        except Exception as e:
            self.add_warning("metadata_error", str(e))
            return DocumentMetadata(file_type=DocumentType.PPTX.value)
    
    def extract_text(self, filepath: str) -> str:
        """Extract all text from PPTX file."""
        text_parts = []
        
        try:
            prs = Presentation(filepath)
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = self._extract_slide_text(slide)
                if slide_text.strip():
                    text_parts.append(f"[Slide {slide_num}]\n{slide_text}")
                    
        except Exception as e:
            self.add_warning("file_error", str(e))
            raise
        
        return '\n\n'.join(text_parts)
    
    def extract_tables(self, filepath: str) -> List[TableData]:
        """Extract all tables from PPTX file."""
        tables = []
        table_index = 0
        
        try:
            prs = Presentation(filepath)
            
            for slide_num, slide in enumerate(prs.slides, 1):
                for shape in slide.shapes:
                    if shape.has_table:
                        try:
                            table = shape.table
                            rows_data = []
                            
                            for row in table.rows:
                                row_data = []
                                for cell in row.cells:
                                    cell_text = cell.text_frame.text.strip() if cell.text_frame else ""
                                    cell_text = re.sub(r'\s+', ' ', cell_text)
                                    row_data.append(cell_text)
                                
                                if any(cell for cell in row_data):
                                    rows_data.append(row_data)
                            
                            if rows_data:
                                headers = rows_data[0] if rows_data else []
                                rows = rows_data[1:] if len(rows_data) > 1 else []
                                
                                tables.append(TableData(
                                    index=table_index,
                                    headers=headers,
                                    rows=rows,
                                    slide=slide_num,
                                ))
                                table_index += 1
                                
                        except Exception as e:
                            self.add_warning(
                                "table_extraction_error",
                                str(e),
                                f"slide {slide_num}"
                            )
                            
        except Exception as e:
            self.add_warning("file_error", str(e))
            raise
        
        return tables
    
    def extract_images(self, filepath: str) -> List[ImageInfo]:
        """Extract image information from PPTX file."""
        images = []
        image_index = 0
        
        try:
            prs = Presentation(filepath)
            
            for slide_num, slide in enumerate(prs.slides, 1):
                for shape in slide.shapes:
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        try:
                            # Get image dimensions from shape
                            images.append(ImageInfo(
                                index=image_index,
                                width=int(shape.width) if shape.width else None,
                                height=int(shape.height) if shape.height else None,
                                slide=slide_num,
                            ))
                            image_index += 1
                            
                        except Exception as e:
                            self.add_warning(
                                "image_extraction_error",
                                str(e),
                                f"slide {slide_num}"
                            )
                            
        except Exception as e:
            self.add_warning("file_error", str(e))
            raise
        
        return images
    
    def extract_slides(self, filepath: str) -> List[SlideContent]:
        """
        Extract content from each slide.
        
        Returns slide-specific information including:
        - Slide title
        - Main content
        - Speaker notes
        - Tables
        - Images
        """
        slides = []
        
        try:
            prs = Presentation(filepath)
            
            for slide_num, slide in enumerate(prs.slides, 1):
                try:
                    # Extract title
                    title = self._extract_slide_title(slide)
                    
                    # Extract main content
                    content = self._extract_slide_text(slide)
                    
                    # Extract speaker notes
                    notes = self._extract_speaker_notes(slide)
                    
                    # Extract tables for this slide
                    slide_tables = self._extract_slide_tables(slide, slide_num)
                    
                    # Extract images for this slide
                    slide_images = self._extract_slide_images(slide, slide_num)
                    
                    slides.append(SlideContent(
                        slide_number=slide_num,
                        title=title,
                        content=content,
                        speaker_notes=notes,
                        tables=slide_tables,
                        images=slide_images,
                    ))
                    
                except Exception as e:
                    self.add_warning(
                        "slide_extraction_error",
                        str(e),
                        f"slide {slide_num}"
                    )
                    
        except Exception as e:
            self.add_warning("file_error", str(e))
            raise
        
        return slides
    
    def _extract_slide_text(self, slide) -> str:
        """Extract all text content from a slide."""
        text_parts = []
        
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text_parts.append(shape.text.strip())
            elif shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    para_text = paragraph.text.strip()
                    if para_text:
                        text_parts.append(para_text)
        
        return '\n'.join(text_parts)
    
    def _extract_slide_title(self, slide) -> Optional[str]:
        """Extract slide title if present."""
        # Check for title placeholder
        if slide.shapes.title:
            return slide.shapes.title.text.strip()
        
        # Look for shapes that might be titles (first text shape)
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                # Check if this is likely a title (short, at top)
                text = shape.text.strip()
                if len(text) < 100:  # Titles are usually short
                    return text
        
        return None
    
    def _extract_speaker_notes(self, slide) -> str:
        """Extract speaker notes from a slide."""
        notes_parts = []
        
        if slide.has_notes_slide:
            notes_slide = slide.notes_slide
            for shape in notes_slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    notes_parts.append(shape.text.strip())
        
        return '\n'.join(notes_parts)
    
    def _extract_slide_tables(self, slide, slide_num: int) -> List[TableData]:
        """Extract tables from a specific slide."""
        tables = []
        local_idx = 0
        
        for shape in slide.shapes:
            if shape.has_table:
                try:
                    table = shape.table
                    rows_data = []
                    
                    for row in table.rows:
                        row_data = []
                        for cell in row.cells:
                            cell_text = cell.text_frame.text.strip() if cell.text_frame else ""
                            row_data.append(cell_text)
                        
                        if any(cell for cell in row_data):
                            rows_data.append(row_data)
                    
                    if rows_data:
                        tables.append(TableData(
                            index=local_idx,
                            headers=rows_data[0],
                            rows=rows_data[1:] if len(rows_data) > 1 else [],
                            slide=slide_num,
                        ))
                        local_idx += 1
                        
                except Exception:
                    pass
        
        return tables
    
    def _extract_slide_images(self, slide, slide_num: int) -> List[ImageInfo]:
        """Extract images from a specific slide."""
        images = []
        local_idx = 0
        
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    images.append(ImageInfo(
                        index=local_idx,
                        width=int(shape.width) if shape.width else None,
                        height=int(shape.height) if shape.height else None,
                        slide=slide_num,
                    ))
                    local_idx += 1
                except Exception:
                    pass
        
        return images
    
    def _create_sections_from_slides(self, slides: List[SlideContent]) -> List[Section]:
        """Create sections based on slide content."""
        sections = []
        
        for slide in slides:
            heading = slide.title or f"Slide {slide.slide_number}"
            
            # Combine content and notes
            content_parts = []
            if slide.content:
                content_parts.append(slide.content)
            if slide.speaker_notes:
                content_parts.append(f"[Speaker Notes]\n{slide.speaker_notes}")
            
            content = '\n\n'.join(content_parts)
            
            sections.append(Section(
                heading=heading,
                content=content,
                level=1,
                slide=slide.slide_number,
                tables=slide.tables,
                images=slide.images,
            ))
        
        return sections
