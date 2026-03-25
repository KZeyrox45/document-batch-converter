"""
Shared fixtures for docbatch tests.

This module provides common test fixtures and utilities used across
all test modules.
"""

import os
import tempfile
from pathlib import Path
from typing import Generator

import pytest
from docx import Document
from pptx import Presentation
from pptx.util import Inches


@pytest.fixture
def temp_dir() -> Generator[Path, None, None]:
    """Create a temporary directory for test files."""
    with tempfile.TemporaryDirectory() as tmpdir:
        yield Path(tmpdir)


@pytest.fixture
def sample_pdf(temp_dir: Path) -> Path:
    """Create a sample PDF file for testing."""
    import pdfplumber
    
    pdf_path = temp_dir / "sample.pdf"
    
    # Create a simple PDF using reportlab (if available) or use pdfplumber
    # For simplicity, we'll create a minimal PDF using the lower-level approach
    from reportlab.lib.pagesizes import letter
    from reportlab.pdfgen import canvas
    
    c = canvas.Canvas(str(pdf_path), pagesize=letter)
    c.drawString(100, 750, "1. Introduction")
    c.drawString(100, 720, "This is the introduction section.")
    c.drawString(100, 690, "It contains some sample text.")
    c.drawString(100, 650, "2. Methodology")
    c.drawString(100, 620, "This is the methodology section.")
    c.drawString(100, 590, "We describe our methods here.")
    c.save()
    
    return pdf_path


@pytest.fixture
def sample_docx(temp_dir: Path) -> Path:
    """Create a sample DOCX file for testing."""
    docx_path = temp_dir / "sample.docx"
    
    doc = Document()
    
    # Add title
    doc.add_heading('Sample Document', 0)
    
    # Add sections
    doc.add_heading('1. Introduction', level=1)
    doc.add_paragraph('This is the introduction section of the document.')
    doc.add_paragraph('It contains multiple paragraphs of text.')
    
    doc.add_heading('1.1 Background', level=2)
    doc.add_paragraph('Background information goes here.')
    
    doc.add_heading('2. Methodology', level=1)
    doc.add_paragraph('This section describes the methodology used.')
    
    # Add a table
    table = doc.add_table(rows=3, cols=2)
    table.cell(0, 0).text = 'Header 1'
    table.cell(0, 1).text = 'Header 2'
    table.cell(1, 0).text = 'Data 1'
    table.cell(1, 1).text = 'Data 2'
    table.cell(2, 0).text = 'Data 3'
    table.cell(2, 1).text = 'Data 4'
    
    doc.save(str(docx_path))
    return docx_path


@pytest.fixture
def sample_pptx(temp_dir: Path) -> Path:
    """Create a sample PPTX file for testing."""
    pptx_path = temp_dir / "sample.pptx"
    
    prs = Presentation()
    
    # Slide 1 - Title slide
    slide_layout = prs.slide_layouts[0]  # Title slide layout
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Presentation Title"
    subtitle.text = "Sample Presentation"
    
    # Add speaker notes
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = "These are speaker notes for slide 1."
    
    # Slide 2 - Content slide
    slide_layout = prs.slide_layouts[1]  # Title and content layout
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Introduction"
    
    # Add content
    body_shape = slide.placeholders[1]
    tf = body_shape.text_frame
    tf.text = "This is the introduction content."
    
    # Slide 3 - Content with table
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    left = Inches(1)
    top = Inches(1)
    width = Inches(4)
    height = Inches(1)
    
    table = slide.shapes.add_table(2, 2, left, top, width, height).table
    table.cell(0, 0).text = "Column A"
    table.cell(0, 1).text = "Column B"
    table.cell(1, 0).text = "Value 1"
    table.cell(1, 1).text = "Value 2"
    
    prs.save(str(pptx_path))
    return pptx_path


@pytest.fixture
def sample_files_dir(temp_dir: Path, sample_pdf: Path, sample_docx: Path, sample_pptx: Path) -> Path:
    """Create a directory with multiple sample files."""
    # Files are already created in temp_dir by fixtures
    return temp_dir


@pytest.fixture
def corrupted_pdf(temp_dir: Path) -> Path:
    """Create a corrupted PDF file for error handling tests."""
    pdf_path = temp_dir / "corrupted.pdf"
    
    # Write invalid PDF content
    with open(pdf_path, 'w') as f:
        f.write("This is not a valid PDF file content.")
    
    return pdf_path


@pytest.fixture
def empty_pdf(temp_dir: Path) -> Path:
    """Create an empty PDF file for edge case tests."""
    pdf_path = temp_dir / "empty.pdf"
    pdf_path.touch()
    return pdf_path


@pytest.fixture
def unsupported_file(temp_dir: Path) -> Path:
    """Create a file with unsupported extension."""
    file_path = temp_dir / "unsupported.xyz"
    file_path.write_text("Some content")
    return file_path
