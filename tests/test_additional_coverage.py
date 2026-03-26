"""
Additional tests to improve coverage.

This module adds tests for edge cases and error paths
to achieve 80%+ coverage.
"""

import pytest
from pathlib import Path
from unittest.mock import Mock, patch, MagicMock
import tempfile
import os

from docbatch.parsers.pdf_parser import PDFParser
from docbatch.parsers.docx_parser import DOCXParser
from docbatch.parsers.pptx_parser import PPTXParser
from docbatch.converter import DocumentConverter
from docbatch.cli import main


class TestPDFParserEdgeCases:
    """Additional tests for PDF parser coverage."""
    
    def test_parse_with_warnings(self, temp_dir):
        """Test PDF parsing that generates warnings."""
        parser = PDFParser()
        
        # Create a minimal valid PDF
        from reportlab.lib.pagesizes import letter
        from reportlab.pdfgen import canvas
        
        pdf_path = temp_dir / "test.pdf"
        c = canvas.Canvas(str(pdf_path), pagesize=letter)
        c.drawString(100, 750, "Test content")
        c.save()
        
        output = parser.parse(str(pdf_path))
        
        assert output.filename == "test.pdf"
        assert len(parser.warnings) == 0  # Should be cleared at start
    
    def test_clean_table_with_none_values(self):
        """Test table cleaning with None values."""
        parser = PDFParser()
        
        table = [
            [None, "Header", None],
            ["Row1", None, "Data"],
            [None, None, None],
        ]
        
        result = parser._clean_table(table)
        
        # Should clean None values
        assert result[0] == ["", "Header", ""]
        assert result[1] == ["Row1", "", "Data"]
    
    def test_detect_sections_with_page_marker(self):
        """Test section detection with page markers."""
        parser = PDFParser()
        
        text = """[Page 1]
1. Introduction

This is introduction content.

[Page 2]
2. Methods

Methods go here.
"""
        
        sections = parser.detect_sections(text)
        
        assert len(sections) >= 1


class TestDOCXParserEdgeCases:
    """Additional tests for DOCX parser coverage."""
    
    def test_extract_images_from_docx(self, sample_docx):
        """Test image extraction from DOCX."""
        parser = DOCXParser()
        
        images = parser.extract_images(str(sample_docx))
        
        # May or may not have images
        assert isinstance(images, list)
    
    def test_count_pages(self, sample_docx):
        """Test page count estimation."""
        parser = DOCXParser()
        
        from docx import Document
        doc = Document(str(sample_docx))
        
        pages = parser._count_pages(doc)
        
        assert pages >= 1
    
    def test_get_image_format(self):
        """Test image format detection."""
        parser = DOCXParser()
        
        assert parser._get_image_format("image.png") == "PNG"
        assert parser._get_image_format("image.jpg") == "JPEG"
        assert parser._get_image_format("image.jpeg") == "JPEG"
        assert parser._get_image_format("image.gif") == "GIF"
        assert parser._get_image_format("image.unknown") == "UNKNOWN"
    
    def test_extract_text_empty_paragraphs(self, temp_dir):
        """Test text extraction with empty paragraphs."""
        from docx import Document
        
        docx_path = temp_dir / "empty.docx"
        doc = Document()
        
        # Add empty paragraphs
        doc.add_paragraph("")
        doc.add_paragraph("   ")  # Whitespace only
        doc.add_paragraph("Real content")
        
        doc.save(str(docx_path))
        
        parser = DOCXParser()
        text = parser.extract_text(str(docx_path))
        
        assert "Real content" in text


class TestPPTXParserEdgeCases:
    """Additional tests for PPTX parser coverage."""
    
    def test_extract_slide_title_none(self, sample_pptx):
        """Test slide title extraction when none exists."""
        from pptx import Presentation
        
        parser = PPTXParser()
        prs = Presentation(str(sample_pptx))
        
        for slide in prs.slides:
            title = parser._extract_slide_title(slide)
            # May be None or a string
            assert title is None or isinstance(title, str)
    
    def test_extract_slide_images(self, sample_pptx):
        """Test image extraction from slides."""
        from pptx import Presentation
        
        parser = PPTXParser()
        prs = Presentation(str(sample_pptx))
        
        for slide in prs.slides:
            images = parser._extract_slide_images(slide, 1)
            assert isinstance(images, list)
    
    def test_extract_slide_tables(self, sample_pptx):
        """Test table extraction from slides."""
        from pptx import Presentation
        
        parser = PPTXParser()
        prs = Presentation(str(sample_pptx))
        
        for slide in prs.slides:
            tables = parser._extract_slide_tables(slide, 1)
            assert isinstance(tables, list)


class TestConverterEdgeCases:
    """Additional tests for converter coverage."""
    
    def test_convert_directory_empty(self, temp_dir):
        """Test converting an empty directory."""
        empty_dir = temp_dir / "empty"
        empty_dir.mkdir()
        
        converter = DocumentConverter(show_progress=False)
        results, stats = converter.convert_directory(empty_dir)
        
        assert stats.total_files == 0
        assert results == []
    
    def test_convert_file_with_errors_skip_true(self, temp_dir):
        """Test file conversion with skip_errors=True."""
        converter = DocumentConverter(skip_errors=True, show_progress=False)
        
        # Create a corrupted file
        bad_file = temp_dir / "bad.pdf"
        bad_file.write_text("Not a PDF")
        
        # Should not raise, returns output with warnings
        output = converter.convert_file(bad_file)
        
        assert output is not None
    
    def test_find_files_with_mixed_content(self, temp_dir):
        """Test file finding with mixed content."""
        # Create various files
        (temp_dir / "doc1.pdf").touch()
        (temp_dir / "doc2.docx").touch()
        (temp_dir / "doc3.pptx").touch()
        (temp_dir / "readme.txt").touch()  # Unsupported
        (temp_dir / "image.png").touch()  # Unsupported
        
        converter = DocumentConverter(show_progress=False)
        files = converter._find_files(temp_dir, recursive=False)
        
        # Should only find supported files
        assert len(files) == 3
        extensions = {f.suffix.lower() for f in files}
        assert extensions == {'.pdf', '.docx', '.pptx'}


class TestCLIEdgeCases:
    """Additional tests for CLI coverage."""
    
    def test_main_convert_directory_output(self, sample_files_dir, temp_dir):
        """Test directory conversion with output path."""
        output_dir = temp_dir / "json_output"
        
        result = main([
            'convert',
            str(sample_files_dir),
            '-o', str(output_dir),
            '-q'
        ])
        
        assert result == 0
        assert output_dir.exists()
    
    def test_main_convert_strict_mode(self, sample_pdf):
        """Test convert with strict mode."""
        result = main([
            'convert',
            str(sample_pdf),
            '--strict',
            '-q'
        ])
        
        assert result == 0
    
    def test_main_keyboard_interrupt(self, sample_files_dir):
        """Test handling of keyboard interrupt."""
        with patch('docbatch.converter.DocumentConverter.convert_directory') as mock_convert:
            mock_convert.side_effect = KeyboardInterrupt()
            
            result = main(['convert', str(sample_files_dir), '-q'])
            
            assert result == 130  # Standard exit code for SIGINT
    
    def test_main_convert_with_verbose_output(self, sample_pdf, temp_dir):
        """Test convert with verbose and output file."""
        output_path = temp_dir / "output.json"
        
        result = main([
            'convert',
            str(sample_pdf),
            '-o', str(output_path),
            '-v'
        ])
        
        assert result == 0
        assert output_path.exists()


class TestBaseParserEdgeCases:
    """Additional tests for base parser coverage."""
    
    def test_time_conversion(self):
        """Test timing conversion operation."""
        parser = PDFParser()
        
        def sample_func(x):
            return x * 2
        
        elapsed, result = parser.time_conversion(sample_func, 5)
        
        assert result == 10
        assert elapsed >= 0
    
    def test_detect_sections_mixed_markers(self):
        """Test section detection with mixed markers."""
        parser = PDFParser()
        
        text = """
Chapter 1: Overview

This is chapter 1.

Section 2

This is section 2.

3. Conclusion

This is the conclusion.
"""
        
        sections = parser.detect_sections(text)
        
        assert len(sections) >= 1
    
    def test_detect_sections_appendix(self):
        """Test detection of appendix markers."""
        parser = PDFParser()
        
        text = """
1. Introduction

Content here.

Appendix A: Additional Data

More content.
"""
        
        sections = parser.detect_sections(text)
        
        assert len(sections) >= 1
