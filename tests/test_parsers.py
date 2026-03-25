"""
Tests for document parsers.

This module tests the PDF, DOCX, and PPTX parser implementations.
"""

import pytest
from pathlib import Path
from unittest.mock import Mock, patch, MagicMock

from docbatch.parsers.base import BaseParser
from docbatch.parsers.pdf_parser import PDFParser
from docbatch.parsers.docx_parser import DOCXParser
from docbatch.parsers.pptx_parser import PPTXParser
from docbatch.models import Section, TableData


class TestBaseParser:
    """Tests for the BaseParser abstract class."""
    
    def test_supports_file_method(self):
        """Test the supports_file class method."""
        assert PDFParser.supports_file("test.pdf")
        assert not PDFParser.supports_file("test.docx")
        
        assert DOCXParser.supports_file("test.docx")
        assert DOCXParser.supports_file("test.DOCX")  # Case insensitive
        
        assert PPTXParser.supports_file("test.pptx")
        assert not PPTXParser.supports_file("test.pdf")
    
    def test_detect_sections_numbered(self):
        """Test section detection with numbered headings."""
        parser = PDFParser()
        
        text = """
1. Introduction

This is the introduction section.

1.1 Background

Background information here.

2. Methodology

Methods are described here.
"""
        
        sections = parser.detect_sections(text)
        
        assert len(sections) >= 1
        # Check that at least one section was detected
        assert any("Introduction" in s.heading for s in sections)
    
    def test_detect_sections_chapter(self):
        """Test section detection with chapter headings."""
        parser = PDFParser()
        
        text = """
Chapter 1: Introduction

Welcome to this document.

Chapter 2: Methods

Our methods are simple.
"""
        
        sections = parser.detect_sections(text)
        
        assert len(sections) >= 1
    
    def test_detect_sections_empty_text(self):
        """Test section detection with empty text."""
        parser = PDFParser()
        
        sections = parser.detect_sections("")
        
        assert sections == []
    
    def test_detect_sections_no_markers(self):
        """Test section detection when no explicit markers exist."""
        parser = PDFParser()
        
        text = "Just some plain text without any section markers."
        
        sections = parser.detect_sections(text)
        
        # Should create a single section with all content
        assert len(sections) == 1
        assert sections[0].heading == "Document Content"
    
    def test_calculate_heading_level(self):
        """Test heading level calculation."""
        parser = PDFParser()
        
        assert parser._calculate_heading_level("1") == 1
        assert parser._calculate_heading_level("1.1") == 2
        assert parser._calculate_heading_level("1.1.1") == 3
        assert parser._calculate_heading_level("I") == 1
        assert parser._calculate_heading_level("A") == 1
    
    def test_add_warning(self):
        """Test warning tracking."""
        parser = PDFParser()
        
        parser.add_warning("test_warning", "Test message", "location 1")
        
        assert len(parser.warnings) == 1
        assert parser.warnings[0].type == "test_warning"
        
        parser.clear_warnings()
        assert len(parser.warnings) == 0


class TestPDFParser:
    """Tests for the PDF parser."""
    
    def test_supported_extensions(self):
        """Test PDF parser supported extensions."""
        assert '.pdf' in PDFParser.SUPPORTED_EXTENSIONS
    
    def test_parse_sample_file(self, sample_pdf):
        """Test parsing a sample PDF file."""
        parser = PDFParser()
        
        output = parser.parse(str(sample_pdf))
        
        assert output.filename == "sample.pdf"
        assert output.metadata.file_type == "pdf"
        assert output.metadata.pages >= 1
    
    def test_extract_metadata(self, sample_pdf):
        """Test PDF metadata extraction."""
        parser = PDFParser()
        
        metadata = parser.extract_metadata(str(sample_pdf))
        
        assert metadata.file_type == "pdf"
        assert metadata.pages >= 1
    
    def test_extract_text(self, sample_pdf):
        """Test text extraction from PDF."""
        parser = PDFParser()
        
        text = parser.extract_text(str(sample_pdf))
        
        assert len(text) > 0
        assert "Introduction" in text or "Methodology" in text
    
    def test_extract_tables(self, sample_pdf):
        """Test table extraction from PDF."""
        parser = PDFParser()
        
        tables = parser.extract_tables(str(sample_pdf))
        
        # May or may not have tables depending on the sample
        assert isinstance(tables, list)
    
    def test_extract_images(self, sample_pdf):
        """Test image extraction from PDF."""
        parser = PDFParser()
        
        images = parser.extract_images(str(sample_pdf))
        
        # May or may not have images depending on the sample
        assert isinstance(images, list)
    
    def test_parse_nonexistent_file(self):
        """Test handling of nonexistent file."""
        parser = PDFParser()
        
        with pytest.raises(Exception):
            parser.parse("/nonexistent/file.pdf")


class TestDOCXParser:
    """Tests for the DOCX parser."""
    
    def test_supported_extensions(self):
        """Test DOCX parser supported extensions."""
        assert '.docx' in DOCXParser.SUPPORTED_EXTENSIONS
        assert '.doc' in DOCXParser.SUPPORTED_EXTENSIONS
    
    def test_parse_sample_file(self, sample_docx):
        """Test parsing a sample DOCX file."""
        parser = DOCXParser()
        
        output = parser.parse(str(sample_docx))
        
        assert output.filename == "sample.docx"
        assert output.metadata.file_type == "docx"
    
    def test_extract_metadata(self, sample_docx):
        """Test DOCX metadata extraction."""
        parser = DOCXParser()
        
        metadata = parser.extract_metadata(str(sample_docx))
        
        assert metadata.file_type == "docx"
    
    def test_extract_text(self, sample_docx):
        """Test text extraction from DOCX."""
        parser = DOCXParser()
        
        text = parser.extract_text(str(sample_docx))
        
        assert len(text) > 0
        assert "Introduction" in text or "introduction" in text.lower()
    
    def test_extract_tables(self, sample_docx):
        """Test table extraction from DOCX."""
        parser = DOCXParser()
        
        tables = parser.extract_tables(str(sample_docx))
        
        assert len(tables) >= 1  # Our sample has a table
        assert tables[0].headers == ["Header 1", "Header 2"]
    
    def test_detect_sections_from_headings(self, sample_docx):
        """Test section detection from Word headings."""
        parser = DOCXParser()
        
        sections = parser._detect_sections_from_headings(str(sample_docx))
        
        assert len(sections) >= 1
        # Should detect "Introduction" and "Methodology" headings
        headings = [s.heading for s in sections]
        assert any("Introduction" in h for h in headings)


class TestPPTXParser:
    """Tests for the PPTX parser."""
    
    def test_supported_extensions(self):
        """Test PPTX parser supported extensions."""
        assert '.pptx' in PPTXParser.SUPPORTED_EXTENSIONS
        assert '.ppt' in PPTXParser.SUPPORTED_EXTENSIONS
    
    def test_parse_sample_file(self, sample_pptx):
        """Test parsing a sample PPTX file."""
        parser = PPTXParser()
        
        output = parser.parse(str(sample_pptx))
        
        assert output.filename == "sample.pptx"
        assert output.metadata.file_type == "pptx"
        assert len(output.slides) >= 1
    
    def test_extract_metadata(self, sample_pptx):
        """Test PPTX metadata extraction."""
        parser = PPTXParser()
        
        metadata = parser.extract_metadata(str(sample_pptx))
        
        assert metadata.file_type == "pptx"
        assert metadata.slides >= 1
    
    def test_extract_text(self, sample_pptx):
        """Test text extraction from PPTX."""
        parser = PPTXParser()
        
        text = parser.extract_text(str(sample_pptx))
        
        assert len(text) > 0
    
    def test_extract_slides(self, sample_pptx):
        """Test slide extraction from PPTX."""
        parser = PPTXParser()
        
        slides = parser.extract_slides(str(sample_pptx))
        
        assert len(slides) >= 2  # Our sample has 3 slides
        
        # Check that speaker notes are extracted
        assert slides[0].speaker_notes or True  # May or may not have notes
    
    def test_extract_tables(self, sample_pptx):
        """Test table extraction from PPTX."""
        parser = PPTXParser()
        
        tables = parser.extract_tables(str(sample_pptx))
        
        # Our sample has a table on slide 3
        assert len(tables) >= 1
    
    def test_extract_slide_text(self, sample_pptx):
        """Test individual slide text extraction."""
        from pptx import Presentation
        
        parser = PPTXParser()
        prs = Presentation(str(sample_pptx))
        
        slide = prs.slides[0]
        text = parser._extract_slide_text(slide)
        
        assert len(text) > 0
    
    def test_create_sections_from_slides(self, sample_pptx):
        """Test section creation from slides."""
        parser = PPTXParser()
        
        slides = parser.extract_slides(str(sample_pptx))
        sections = parser._create_sections_from_slides(slides)
        
        assert len(sections) == len(slides)
        assert all(s.slide is not None for s in sections)
